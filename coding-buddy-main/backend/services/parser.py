import fitz  # PyMuPDF
from pptx import Presentation
import pytesseract
from PIL import Image

def parse_file(file_path: str, file_type: str) -> list[dict]:
    extracted_data = []
    
    file_type = file_type.lower()
    
    if file_type == "pdf":
        doc = fitz.open(file_path)
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text()
            if text.strip() != "":
                extracted_data.append({
                    "text": text.strip(),
                    "page": i + 1
                })
        doc.close()
        
    elif file_type == "pptx":
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text.strip())
            
            combined_text = "\n".join(slide_text)
            if combined_text.strip() != "":
                extracted_data.append({
                    "text": combined_text.strip(),
                    "page": i + 1
                })
                
    elif file_type in ["png", "jpg", "jpeg"]:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        if text.strip() != "":
            extracted_data.append({
                "text": text.strip(),
                "page": 1
            })
            
    return extracted_data
